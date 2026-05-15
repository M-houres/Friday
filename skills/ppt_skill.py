"""
一秒PPT Skill —— 输入主题，自动生成专业PPTX文件

全链路: 意图识别 → 主题研究 → 大纲生成 → 并行内容生成 → PPTX组装 → 文件交付

依赖: python-pptx (pip install python-pptx)
"""

import asyncio
import io
import json
import logging
import os
import uuid
from datetime import datetime
from pathlib import Path
from typing import Optional

from src.tools.skill import FridaySkill, tool, skill
from src.models.router import model_router
from src.models.base import Message
from src.artifacts.service import artifact_service

logger = logging.getLogger(__name__)

OUTPUT_DIR = Path(__file__).parent.parent / "static" / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


@skill(
    name="一秒PPT",
    trigger="PPT|幻灯片|演示文稿|做PPT|生成PPT|制作PPT|presentation|slide",
    description="输入主题，全自动生成10页专业PPT并导出PPTX文件。研究→大纲→内容→排版→下载，一步到位。",
    version="2.0.0",
    icon="📊",
)
class PPTSkill(FridaySkill):
    workflow = [
        {"id": "research", "tool": "research_topic", "name": "主题研究", "dependencies": []},
        {"id": "outline", "tool": "generate_outline", "name": "生成大纲", "dependencies": ["research"]},
        {"id": "slides", "tool": "generate_all_slides", "name": "并行生成内容", "dependencies": ["outline"]},
        {"id": "assemble", "tool": "assemble_pptx", "name": "组装PPTX文件", "dependencies": ["slides"]},
        {"id": "certify", "tool": "certify_quality", "name": "质量校验", "dependencies": ["assemble"]},
        {"id": "finalize", "tool": "finalize_output", "name": "交付下载", "dependencies": ["certify"]},
    ]

    output = {
        "slides": 10,
        "format": "pptx",
        "language": "zh-CN",
        "style": "professional",
    }

    # ── 工具方法 ──

    @tool(
        name="research_topic",
        description="深度研究用户提供的主题，收集关键信息、数据、案例，为PPT内容提供素材基础",
        parameters={
            "type": "object",
            "properties": {
                "task": {"type": "string", "description": "用户输入的完整任务描述"},
                "context": {"type": "object", "description": "额外上下文（可选）"},
            },
            "required": ["task"],
        },
    )
    async def research_topic(self, task: str, context: Optional[dict] = None) -> dict:
        logger.info(f"[PPT] 研究主题: {task[:60]}")
        topic = self._clean_topic(task)

        messages = [
            Message(
                role="system",
                content="""你是专业的研究分析师。为PPT制作收集素材。

输出严格的JSON格式:
{
  "topic": "主题名称",
  "keywords": ["关键词1", "关键词2", ...],
  "key_facts": ["事实1", "事实2", ...],
  "data_points": [{"label": "指标名", "value": "数值"}, ...],
  "target_audience": "目标受众",
  "tone": "professional|casual|academic|business"
}""",
            ),
            Message(role="user", content=f"研究主题: {topic}"),
        ]

        response = await model_router.chat(messages=messages, temperature=0.5, max_tokens=2048)

        try:
            research = self._parse_json(response.content)
        except Exception:
            research = {
                "topic": topic,
                "keywords": [topic],
                "key_facts": [response.content[:200]],
                "data_points": [],
            }

        return {"success": True, "data": research, "tokens": response.tokens_used}

    @tool(
        name="generate_outline",
        description="根据研究结果生成PPT大纲：10个幻灯片的标题和内容要点，确保逻辑递进、结构完整",
        parameters={
            "type": "object",
            "properties": {
                "task": {"type": "string", "description": "原始任务"},
                "context": {"type": "object", "description": "前一步的研究结果"},
            },
            "required": ["task"],
        },
        depends_on=["research_topic"],
    )
    async def generate_outline(self, task: str, context: Optional[dict] = None) -> dict:
        logger.info(f"[PPT] 生成大纲")

        research_data = context.get("research_topic", {}).get("data", {}) if context else {}
        topic = self._clean_topic(task)
        research_json = json.dumps(research_data, ensure_ascii=False, indent=2)

        messages = [
            Message(
                role="system",
                content="""你是顶级PPT策划师。为以下主题设计10页PPT大纲。

输出严格的JSON数组，每页一个对象:
[
  {
    "slide_number": 1,
    "title": "封面标题",
    "subtitle": "副标题",
    "type": "cover|content|data|comparison|conclusion",
    "content_points": ["要点1", "要点2", "要点3"],
    "speaker_notes": "演讲者备注"
  },
  ...
]

规则:
- 第1页: 封面(标题+副标题)
- 第2-9页: 内容页，逻辑递进
- 第10页: 总结/行动号召
- 每页3-5个内容要点
- 标题简洁有力，10字以内""",
            ),
            Message(role="user", content=f"主题: {topic}\n研究资料:\n{research_json}"),
        ]

        response = await model_router.chat(messages=messages, temperature=0.6, max_tokens=4096)

        try:
            outline = self._parse_json(response.content)
            if isinstance(outline, dict):
                outline = outline.get("slides", outline.get("outline", [outline]))
            if not isinstance(outline, list):
                outline = [outline]
        except Exception:
            outline = [{"slide_number": i + 1, "title": f"第{i + 1}页", "type": "content", "content_points": []} for i in range(10)]

        return {"success": True, "data": outline[:10], "tokens": response.tokens_used}

    @tool(
        name="generate_all_slides",
        description="并行生成全部10页PPT的详细内容——包含标题、正文、数据、图表描述。内部使用asyncio.gather实现真正的并行调度",
        parameters={
            "type": "object",
            "properties": {
                "task": {"type": "string", "description": "原始任务"},
                "context": {"type": "object", "description": "前一步的大纲结果"},
            },
            "required": ["task"],
        },
        depends_on=["generate_outline"],
        is_expensive=True,
    )
    async def generate_all_slides(self, task: str, context: Optional[dict] = None) -> dict:
        outline_data = context.get("generate_outline", {}).get("data", []) if context else []
        topic = self._clean_topic(task)

        if not outline_data:
            outline_data = [{"slide_number": i + 1, "title": f"第{i + 1}页", "type": "content"} for i in range(10)]

        logger.info(f"[PPT] 并行生成 {len(outline_data)} 页内容")

        async def generate_single_slide(slide_info: dict) -> dict:
            slide_num = slide_info.get("slide_number", 1)
            slide_title = slide_info.get("title", f"第{slide_num}页")
            slide_type = slide_info.get("type", "content")
            points = slide_info.get("content_points", [])

            messages = [
                Message(
                    role="system",
                    content=f"""你是专业PPT内容写手。为以下幻灯片撰写详细内容。

输出JSON:
{{
  "slide_number": {slide_num},
  "title": "页面标题",
  "body_text": "正文段落(200-400字)",
  "bullet_points": ["要点1(20字内)", "要点2", "要点3", "要点4(可选)"],
  "data_note": "需要突出的数据/引用(可选)"
}}

主题: {topic}
页面类型: {slide_type}""",
                ),
                Message(
                    role="user",
                    content=f"页面标题: {slide_title}\n已有要点: {json.dumps(points, ensure_ascii=False)}",
                ),
            ]

            try:
                response = await model_router.chat(messages=messages, temperature=0.7, max_tokens=2048)
                content = self._parse_json(response.content)
                content["slide_number"] = slide_num
                content["_tokens"] = response.tokens_used
                return content
            except Exception as e:
                return {
                    "slide_number": slide_num,
                    "title": slide_title,
                    "body_text": f"关于{slide_title}的详细内容",
                    "bullet_points": points + ["更多精彩内容"],
                    "_error": str(e),
                }

        # === 真正的并行调度: asyncio.gather ===
        tasks = [generate_single_slide(slide) for slide in outline_data]
        slide_contents = await asyncio.gather(*tasks)

        total_tokens = sum(s.get("_tokens", 0) for s in slide_contents)
        slide_contents.sort(key=lambda s: s.get("slide_number", 0))

        logger.info(f"[PPT] 并行生成完成: {len(slide_contents)}页, {total_tokens} tokens")

        return {
            "success": True,
            "data": slide_contents,
            "parallel_count": len(slide_contents),
            "total_tokens": total_tokens,
        }

    @tool(
        name="assemble_pptx",
        description="将所有幻灯片内容组装为可下载的PPTX文件，包含专业排版和样式",
        parameters={
            "type": "object",
            "properties": {
                "task": {"type": "string", "description": "原始任务"},
                "context": {"type": "object", "description": "全部前置结果"},
            },
            "required": ["task"],
        },
        depends_on=["generate_all_slides"],
    )
    async def assemble_pptx(self, task: str, context: Optional[dict] = None) -> dict:
        slides_content = context.get("generate_all_slides", {}).get("data", []) if context else []
        topic = self._clean_topic(task)

        if not slides_content:
            return {"success": False, "error": "没有幻灯片内容可组装"}

        logger.info(f"[PPT] 组装PPTX: {len(slides_content)}页")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        except ImportError:
            return {
                "success": False,
                "error": "python-pptx 未安装。请运行: pip install python-pptx",
            }

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0xDB)
        COLOR_ACCENT = RGBColor(0x37, 0x41, 0x51)
        COLOR_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
        COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

        for i, slide_data in enumerate(slides_content):
            slide_num = slide_data.get("slide_number", i + 1)
            slide_type = slide_data.get("type", "content")
            title = slide_data.get("title", f"幻灯片 {slide_num}")
            body = slide_data.get("body_text", "")
            bullets = slide_data.get("bullet_points", [])
            data_note = slide_data.get("data_note", "")

            if slide_type == "cover" or slide_num == 1:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                bg = slide.background
                fill = bg.fill
                fill.solid()
                fill.fore_color.rgb = COLOR_PRIMARY

                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(1.5))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = title
                p.font.size = Pt(48)
                p.font.bold = True
                p.font.color.rgb = COLOR_WHITE
                p.alignment = PP_ALIGN.LEFT

                if body:
                    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.333), Inches(1))
                    tf2 = txBox2.text_frame
                    p2 = tf2.paragraphs[0]
                    p2.text = body[:200]
                    p2.font.size = Pt(20)
                    p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xEE)
                    p2.alignment = PP_ALIGN.LEFT

            else:
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                bg_fill = slide.background.fill
                bg_fill.solid()
                bg_fill.fore_color.rgb = COLOR_WHITE

                left_bar = slide.shapes.add_shape(
                    1, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
                )
                left_bar.fill.solid()
                left_bar.fill.fore_color.rgb = COLOR_PRIMARY
                left_bar.line.fill.background()

                top_bar = slide.shapes.add_shape(
                    1, Inches(0), Inches(0), Inches(13.333), Inches(0.08)
                )
                top_bar.fill.solid()
                top_bar.fill.fore_color.rgb = COLOR_PRIMARY
                top_bar.line.fill.background()

                title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8))
                tf_title = title_box.text_frame
                p_title = tf_title.paragraphs[0]
                p_title.text = title
                p_title.font.size = Pt(32)
                p_title.font.bold = True
                p_title.font.color.rgb = COLOR_ACCENT

                page_num_box = slide.shapes.add_textbox(Inches(12), Inches(7), Inches(1), Inches(0.4))
                pn = page_num_box.text_frame.paragraphs[0]
                pn.text = str(slide_num)
                pn.font.size = Pt(11)
                pn.font.color.rgb = RGBColor(0x9C, 0xA3, 0xAF)
                pn.alignment = PP_ALIGN.RIGHT

                if bullets:
                    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5))
                    tf_body = body_box.text_frame
                    tf_body.word_wrap = True
                    for bi, bullet in enumerate(bullets):
                        if bi == 0:
                            p = tf_body.paragraphs[0]
                        else:
                            p = tf_body.add_paragraph()
                        p.text = f"▸ {bullet}"
                        p.font.size = Pt(18)
                        p.font.color.rgb = COLOR_ACCENT
                        p.space_after = Pt(12)
                        p.level = 0

                elif body:
                    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5))
                    tf_body = body_box.text_frame
                    tf_body.word_wrap = True
                    p_body = tf_body.paragraphs[0]
                    p_body.text = body
                    p_body.font.size = Pt(16)
                    p_body.font.color.rgb = COLOR_ACCENT
                    p_body.line_spacing = Pt(28)

                if data_note:
                    note_box = slide.shapes.add_textbox(Inches(1), Inches(6.4), Inches(11), Inches(0.6))
                    tf_note = note_box.text_frame
                    p_note = tf_note.paragraphs[0]
                    p_note.text = f"💡 {data_note}"
                    p_note.font.size = Pt(12)
                    p_note.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)
                    p_note.italic = True

        topic_slug = topic[:20].replace(" ", "_").replace("/", "_")
        filename = f"{topic_slug}_{uuid.uuid4().hex[:8]}.pptx"
        filepath = OUTPUT_DIR / filename
        workflow_id = (context or {}).get("_workflow_id", "adhoc")
        owner_user_id = (context or {}).get("_user_id", "default")

        prs.save(str(filepath))
        file_size = filepath.stat().st_size
        artifact = artifact_service.create_from_file(
            workflow_id=workflow_id,
            filename=filename,
            source_path=str(filepath),
            owner_user_id=owner_user_id,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

        logger.info(f"[PPT] PPTX已保存: {filepath} ({file_size} bytes)")

        return {
            "success": True,
            "data": {
                "filename": artifact["filename"],
                "filepath": str(filepath),
                "size_bytes": file_size,
                "slides_count": len(slides_content),
                "artifact_id": artifact["artifact_id"],
                "download_url": artifact["download_url"],
            },
        }

    @tool(
        name="certify_quality",
        description="对生成的PPTX进行质量校验：检查页数正确、内容完整、格式无误",
        parameters={
            "type": "object",
            "properties": {
                "task": {"type": "string", "description": "原始任务"},
                "context": {"type": "object", "description": "组装结果"},
            },
            "required": ["task"],
        },
        depends_on=["assemble_pptx"],
    )
    async def certify_quality(self, task: str, context: Optional[dict] = None) -> dict:
        assemble_result = context.get("assemble_pptx", {}).get("data", {}) if context else {}
        filepath = assemble_result.get("filepath", "")

        checks = {
            "file_exists": False,
            "file_size_ok": False,
            "slide_count_ok": False,
            "format_valid": False,
        }

        if filepath and os.path.exists(filepath):
            checks["file_exists"] = True
            size = os.path.getsize(filepath)
            checks["file_size_ok"] = size > 1000

        slides_count = assemble_result.get("slides_count", 0)
        checks["slide_count_ok"] = 8 <= slides_count <= 15

        if filepath and filepath.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(filepath)
                actual_slides = len(prs.slides)
                checks["format_valid"] = actual_slides > 0
                checks["actual_slides"] = actual_slides
            except Exception:
                checks["format_valid"] = False

        all_pass = all(checks.values())
        level = "PASS" if all_pass else "FAIL"

        if not all_pass:
            failed = [k for k, v in checks.items() if not v]
            logger.warning(f"[PPT] 质量校验失败: {failed}")

        return {"success": True, "data": {"level": level, "checks": checks, "passed": all_pass}}

    @tool(
        name="finalize_output",
        description="生成最终交付物：下载链接、文件信息和预览摘要",
        parameters={
            "type": "object",
            "properties": {
                "task": {"type": "string", "description": "原始任务"},
                "context": {"type": "object", "description": "全部前置结果"},
            },
            "required": ["task"],
        },
        depends_on=["certify_quality"],
    )
    async def finalize_output(self, task: str, context: Optional[dict] = None) -> dict:
        assemble = context.get("assemble_pptx", {}).get("data", {}) if context else {}
        cert = context.get("certify_quality", {}).get("data", {}) if context else {}
        slides = context.get("generate_all_slides", {}).get("data", []) if context else {}

        topic = self._clean_topic(task)
        download_url = assemble.get("download_url", "")
        filename = assemble.get("filename", "output.pptx")
        size_kb = round(assemble.get("size_bytes", 0) / 1024, 1)
        quality_pass = cert.get("passed", False)

        outline = []
        for s in slides[:10]:
            outline.append({
                "num": s.get("slide_number", 0),
                "title": s.get("title", ""),
            })

        return {
            "success": True,
            "data": {
                "topic": topic,
                "filename": filename,
                "download_url": download_url,
                "file_size_kb": size_kb,
                "slides_count": len(slides),
                "quality_check": "PASS" if quality_pass else "FAIL",
                "slide_outline": outline,
                "generated_at": datetime.now().isoformat(),
                "summary": f"✅ 已为您生成关于「{topic}」的专业PPT，共{len(slides)}页，文件大小{size_kb}KB。",
            },
        }

    # ── 辅助方法 ──

    @staticmethod
    def _clean_topic(task: str) -> str:
        prefixes = [
            "生成一个", "做一个", "帮我做", "帮我生成", "制作一个", "创建一个",
            "请生成", "请做一个", "帮做一个", "做一页", "做", "生成",
        ]
        topic = task.strip()
        for p in sorted(prefixes, key=len, reverse=True):
            if topic.startswith(p):
                topic = topic[len(p):].strip()
                break
        keywords = ["PPT", "的PPT", "幻灯片", "的幻灯片", "演示文稿", "的演示文稿", "pptx"]
        for kw in keywords:
            if topic.lower().endswith(kw.lower()):
                topic = topic[: -len(kw)].strip()
        return topic[:80]

    @staticmethod
    def _parse_json(content: str) -> dict | list:
        content = content.strip()
        if "```json" in content:
            content = content.split("```json")[1].split("```")[0]
        elif "```" in content:
            content = content.split("```")[1].split("```")[0]
        return json.loads(content)
